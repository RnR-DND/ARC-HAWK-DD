"""
PowerPoint (PPTX) connector for ARC-HAWK-DD scanner.
Requires: python-pptx
Scans text content in slides, speaker notes, and shapes.
"""

import os
from pathlib import Path

from hawk_scanner.internals import system
from hawk_scanner.internals.validation_integration import validate_findings
from rich.console import Console

console = Console()


def scan_pptx_file(args, file_path, profile_name):
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        system.print_error(args, "python-pptx not installed. Run: pip install python-pptx")
        return []

    results = []
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            texts_to_scan = []

            # All shapes (text frames, tables)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts_to_scan.append(('shape', text))
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                texts_to_scan.append(('table_cell', text))

            # Speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    texts_to_scan.append(('speaker_notes', notes_text))

            for location, text in texts_to_scan:
                matches = system.match_strings(args, text)
                if matches:
                    validated = validate_findings(matches, args)
                    if validated:
                        for match in validated:
                            results.append({
                                'host': str(Path(file_path).parent),
                                'file_path': file_path,
                                'slide': slide_num,
                                'location': location,
                                'pattern_name': match['pattern_name'],
                                'matches': match['matches'],
                                'sample_text': match['sample_text'],
                                'profile': profile_name,
                                'data_source': 'pptx',
                            })
    except Exception as e:
        system.print_error(args, f"Error reading PPTX {file_path}: {e}")

    return results


def execute(args):
    results = []
    system.print_info(args, "Running checks for PPTX sources")
    connections = system.get_connection(args)

    sources_config = connections.get('sources', {})
    pptx_config = sources_config.get('pptx')
    if not pptx_config:
        system.print_error(args, "No PPTX connection details found in connection.yml")
        return results

    for key, config in pptx_config.items():
        paths = config.get('paths', [])
        recursive = config.get('recursive', True)

        for base_path in paths:
            if not os.path.exists(base_path):
                system.print_error(args, f"Path does not exist: {base_path}")
                continue

            if os.path.isfile(base_path) and base_path.lower().endswith('.pptx'):
                results += scan_pptx_file(args, base_path, key)
            elif os.path.isdir(base_path):
                pattern = '**/*.pptx' if recursive else '*.pptx'
                for pptx_file in Path(base_path).glob(pattern):
                    system.print_info(args, f"Scanning: {pptx_file}")
                    results += scan_pptx_file(args, str(pptx_file), key)

    return results
